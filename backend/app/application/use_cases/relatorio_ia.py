"""Relatório Executivo IA (Fase 8 — V4).

Gera relatórios executivos em PDF, Word e PowerPoint a partir do diagnóstico
IA, score, oportunidades e benchmark setorial. Conteúdo:
Resumo Executivo, KPIs, Benchmark, Economia Potencial, Oportunidades,
Plano de Ação, Conclusões.

A IA produz a narrativa (no diagnóstico); este módulo monta os documentos.
Paleta GD Conecta: Indigo #2D3561, Amber #C9A84C, Blue #0077A8, Ivory #F5F1EB.
"""
from __future__ import annotations

import logging
from datetime import datetime
from io import BytesIO

from sqlalchemy.orm import Session

from app.application.use_cases.diagnostico_ia import DiagnosticoIAUseCase
from app.application.use_cases.score_logistico import ScoreLogisticoUseCase
from app.application.use_cases.oportunidades import OportunidadesUseCase
from app.application.use_cases.benchmark_setorial import BenchmarkSetorialUseCase
from app.infrastructure.database.models import EmpresaModel

logger = logging.getLogger(__name__)

INDIGO = "#2D3561"
AMBER = "#C9A84C"
BLUE = "#0077A8"
IVORY = "#F5F1EB"


def _br(valor: float, casas: int = 2) -> str:
    s = f"{(valor or 0):,.{casas}f}"
    return s.replace(",", "X").replace(".", ",").replace("X", ".")


def _moeda(v: float) -> str:
    return f"R$ {_br(v)}"


class RelatorioIAUseCase:
    def __init__(self, db: Session):
        self.db = db

    # ─────────── Coleta de dados ───────────
    def _coletar(self, empresa_id: int) -> dict:
        empresa = self.db.query(EmpresaModel).filter(EmpresaModel.id == empresa_id).first()
        nome_empresa = empresa.razao_social if empresa else f"Empresa {empresa_id}"

        diag_uc = DiagnosticoIAUseCase(self.db)
        diag = diag_uc.ultimo(empresa_id) or diag_uc.gerar(empresa_id)

        score = ScoreLogisticoUseCase(self.db).ultimo(empresa_id) or \
            ScoreLogisticoUseCase(self.db).calcular(empresa_id)

        oportunidades = OportunidadesUseCase(self.db).listar(empresa_id)
        setorial = BenchmarkSetorialUseCase(self.db).comparar(empresa_id)

        return {
            "empresa": nome_empresa,
            "data": datetime.now().strftime("%d/%m/%Y"),
            "diag": diag,
            "score": score,
            "oportunidades": oportunidades,
            "setorial": setorial,
        }

    def gerar(self, empresa_id: int, formato: str = "pdf") -> bytes:
        """Gera o relatório no formato solicitado. Retorna os bytes do arquivo."""
        dados = self._coletar(empresa_id)
        formato = formato.lower()
        if formato == "pdf":
            return self._gerar_pdf(dados)
        if formato in ("word", "docx"):
            return self._gerar_word(dados)
        if formato in ("powerpoint", "pptx", "ppt"):
            return self._gerar_pptx(dados)
        raise ValueError(f"Formato '{formato}' não suportado (use pdf, word ou powerpoint).")

    # ═══════════════════ PDF ═══════════════════
    def _gerar_pdf(self, d: dict) -> bytes:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
        )

        buf = BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
        base = getSampleStyleSheet()
        base.add(ParagraphStyle("T", parent=base["Title"], textColor=colors.HexColor(INDIGO), fontSize=22, spaceAfter=2))
        base.add(ParagraphStyle("Sub", parent=base["Normal"], textColor=colors.HexColor(BLUE),
                                fontSize=14, spaceBefore=16, spaceAfter=6, fontName="Helvetica-Bold"))
        base.add(ParagraphStyle("Corpo", parent=base["Normal"], fontSize=10, leading=15, spaceAfter=6))
        base.add(ParagraphStyle("Info", parent=base["Normal"], fontSize=10, textColor=colors.HexColor("#666666")))

        diag, score = d["diag"], d["score"]
        el = []
        el.append(Paragraph("Relatório Executivo de Inteligência Logística", base["T"]))
        el.append(Paragraph(f"{d['empresa']} · {d['data']}", base["Info"]))
        el.append(Spacer(1, 0.4 * cm))

        # KPIs
        el.append(Paragraph("Indicadores-Chave", base["Sub"]))
        kpi = [
            ["Score Logístico", f"{score.score_total}/100 ({score.classificacao})"],
            ["Economia Potencial", f"{_moeda(diag.economia_potencial)}/ano"],
            ["Oportunidades Abertas", str(len([o for o in d["oportunidades"] if o.status == "ABERTA"]))],
        ]
        t = Table(kpi, colWidths=[7 * cm, 9 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor(IVORY)),
            ("TEXTCOLOR", (0, 0), (0, -1), colors.HexColor(INDIGO)),
            ("FONTNAME", (0, 0), (0, -1), "Helvetica-Bold"),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#DDDDDD")),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("TOPPADDING", (0, 0), (-1, -1), 6), ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        el.append(t)

        # Seções narrativas
        secoes = [
            ("Resumo Executivo", diag.resumo_executivo),
            ("Pontos Fortes", diag.pontos_fortes),
            ("Pontos de Atenção", diag.pontos_atencao),
            ("Principais Desvios", diag.principais_desvios),
            ("Plano de Ação", diag.plano_acao),
        ]
        for titulo, conteudo in secoes:
            if conteudo:
                el.append(Paragraph(titulo, base["Sub"]))
                for linha in conteudo.split("\n"):
                    if linha.strip():
                        el.append(Paragraph(linha.strip(), base["Corpo"]))

        # Oportunidades
        if d["oportunidades"]:
            el.append(Paragraph("Oportunidades Identificadas", base["Sub"]))
            linhas = [["Tipo", "Título", "Economia/ano", "Impacto"]]
            for o in d["oportunidades"][:8]:
                linhas.append([o.tipo, o.titulo[:40],
                               _moeda(o.economia_estimada) if o.economia_estimada else "—", o.impacto])
            to = Table(linhas, colWidths=[3 * cm, 7 * cm, 3.5 * cm, 2.5 * cm])
            to.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(INDIGO)),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(IVORY)]),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#DDDDDD")),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
            ]))
            el.append(to)

        el.append(Spacer(1, 0.5 * cm))
        el.append(Paragraph("Conclusões", base["Sub"]))
        el.append(Paragraph(
            f"A operação apresenta score {score.score_total}/100 ({score.classificacao}), "
            f"com economia potencial estimada de {_moeda(diag.economia_potencial)}/ano. "
            "Recomenda-se priorizar as oportunidades de maior impacto e menor complexidade "
            "para captura rápida de resultado.", base["Corpo"]))

        # Rodapé com aviso de modo simulado
        if diag.modelo_usado and "simulado" in diag.modelo_usado:
            el.append(Spacer(1, 0.3 * cm))
            el.append(Paragraph(
                "<i>Relatório gerado em modo simulado. Os números são reais (calculados "
                "via SQL); a narrativa é demonstrativa.</i>", base["Info"]))

        doc.build(el)
        return buf.getvalue()

    # ═══════════════════ WORD ═══════════════════
    def _gerar_word(self, d: dict) -> bytes:
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        diag, score = d["diag"], d["score"]
        doc = Document()

        titulo = doc.add_heading("Relatório Executivo de Inteligência Logística", level=0)
        titulo.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph(f"{d['empresa']} · {d['data']}")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_heading("Indicadores-Chave", level=1)
        tbl = doc.add_table(rows=3, cols=2)
        tbl.style = "Light Grid Accent 1"
        kpis = [
            ("Score Logístico", f"{score.score_total}/100 ({score.classificacao})"),
            ("Economia Potencial", f"{_moeda(diag.economia_potencial)}/ano"),
            ("Oportunidades Abertas", str(len([o for o in d["oportunidades"] if o.status == "ABERTA"]))),
        ]
        for i, (k, v) in enumerate(kpis):
            tbl.rows[i].cells[0].text = k
            tbl.rows[i].cells[1].text = v

        for titulo_s, conteudo in [
            ("Resumo Executivo", diag.resumo_executivo),
            ("Pontos Fortes", diag.pontos_fortes),
            ("Pontos de Atenção", diag.pontos_atencao),
            ("Principais Desvios", diag.principais_desvios),
            ("Plano de Ação", diag.plano_acao),
        ]:
            if conteudo:
                doc.add_heading(titulo_s, level=1)
                for linha in conteudo.split("\n"):
                    if linha.strip():
                        doc.add_paragraph(linha.strip())

        if d["oportunidades"]:
            doc.add_heading("Oportunidades Identificadas", level=1)
            t = doc.add_table(rows=1, cols=4)
            t.style = "Light Grid Accent 1"
            hdr = t.rows[0].cells
            hdr[0].text, hdr[1].text, hdr[2].text, hdr[3].text = "Tipo", "Título", "Economia/ano", "Impacto"
            for o in d["oportunidades"][:8]:
                c = t.add_row().cells
                c[0].text = o.tipo
                c[1].text = o.titulo[:50]
                c[2].text = _moeda(o.economia_estimada) if o.economia_estimada else "—"
                c[3].text = o.impacto

        doc.add_heading("Conclusões", level=1)
        doc.add_paragraph(
            f"A operação apresenta score {score.score_total}/100 ({score.classificacao}), "
            f"com economia potencial estimada de {_moeda(diag.economia_potencial)}/ano. "
            "Recomenda-se priorizar as oportunidades de maior impacto e menor complexidade.")

        if diag.modelo_usado and "simulado" in diag.modelo_usado:
            p = doc.add_paragraph()
            run = p.add_run("Relatório gerado em modo simulado. Os números são reais "
                            "(calculados via SQL); a narrativa é demonstrativa.")
            run.italic = True
            run.font.size = Pt(8)

        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()

    # ═══════════════════ POWERPOINT ═══════════════════
    def _gerar_pptx(self, d: dict) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        diag, score = d["diag"], d["score"]
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        indigo = RGBColor(0x2D, 0x35, 0x61)
        amber = RGBColor(0xC9, 0xA8, 0x4C)
        blue = RGBColor(0x00, 0x77, 0xA8)
        branco = RGBColor(0xFF, 0xFF, 0xFF)

        def _fundo(slide, cor):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = cor

        def _caixa(slide, texto, left, top, width, height, tam=18, cor=None, bold=False, align=PP_ALIGN.LEFT):
            tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = texto
            run.font.size = Pt(tam)
            run.font.bold = bold
            run.font.color.rgb = cor or RGBColor(0x33, 0x33, 0x33)
            return tb

        # Slide 1 — Capa
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _fundo(s, indigo)
        _caixa(s, "Relatório Executivo", 1, 2.3, 11, 1, tam=44, cor=branco, bold=True, align=PP_ALIGN.CENTER)
        _caixa(s, "Inteligência Logística", 1, 3.4, 11, 1, tam=30, cor=amber, bold=True, align=PP_ALIGN.CENTER)
        _caixa(s, f"{d['empresa']} · {d['data']}", 1, 4.6, 11, 0.6, tam=18, cor=branco, align=PP_ALIGN.CENTER)

        # Slide 2 — KPIs
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _fundo(s, branco)
        _caixa(s, "Indicadores-Chave", 0.7, 0.4, 12, 0.8, tam=30, cor=indigo, bold=True)
        kpis = [
            ("Score Logístico", f"{score.score_total}/100", score.classificacao, indigo),
            ("Economia Potencial", f"{_moeda(diag.economia_potencial)}", "por ano", blue),
            ("Oportunidades", str(len([o for o in d["oportunidades"] if o.status == "ABERTA"])), "abertas", amber),
        ]
        for i, (rotulo, valor, leg, cor) in enumerate(kpis):
            left = 0.7 + i * 4.1
            card = s.shapes.add_shape(1, Inches(left), Inches(1.8), Inches(3.8), Inches(2.6))
            card.fill.solid(); card.fill.fore_color.rgb = cor
            card.line.color.rgb = cor
            _caixa(s, rotulo, left + 0.2, 2.1, 3.4, 0.6, tam=15, cor=branco, align=PP_ALIGN.CENTER)
            _caixa(s, valor, left + 0.2, 2.8, 3.4, 1, tam=34, cor=branco, bold=True, align=PP_ALIGN.CENTER)
            _caixa(s, leg, left + 0.2, 3.8, 3.4, 0.5, tam=13, cor=branco, align=PP_ALIGN.CENTER)

        # Slide 3 — Resumo Executivo
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _fundo(s, branco)
        _caixa(s, "Resumo Executivo", 0.7, 0.4, 12, 0.8, tam=30, cor=indigo, bold=True)
        _caixa(s, (diag.resumo_executivo or "")[:900], 0.7, 1.5, 12, 5.5, tam=16)

        # Slide 4 — Oportunidades
        if d["oportunidades"]:
            s = prs.slides.add_slide(prs.slide_layouts[6])
            _fundo(s, branco)
            _caixa(s, "Oportunidades Identificadas", 0.7, 0.4, 12, 0.8, tam=30, cor=indigo, bold=True)
            y = 1.6
            for o in d["oportunidades"][:6]:
                eco = f" · {_moeda(o.economia_estimada)}/ano" if o.economia_estimada else ""
                _caixa(s, f"• [{o.tipo}] {o.titulo}{eco}", 0.9, y, 11.5, 0.6, tam=15, cor=indigo)
                y += 0.85

        # Slide 5 — Plano de Ação / Conclusões
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _fundo(s, indigo)
        _caixa(s, "Plano de Ação", 0.7, 0.4, 12, 0.8, tam=30, cor=amber, bold=True)
        _caixa(s, (diag.plano_acao or "Priorizar oportunidades de maior impacto.")[:800],
               0.7, 1.5, 12, 5, tam=16, cor=branco)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
